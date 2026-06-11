#!/usr/bin/env python3
"""Extract speaker notes from a .pptx (one line per slide, numbered format).

Output format matches scripts/utils/portfolio-notes.js numbered mode:
  1: first slide notes (whitespace normalized)
  2:
  3: third slide notes

Requires: pnpm run setup-deck-python (project venv with python-pptx)
"""

from __future__ import annotations

import argparse
import sys

try:
    from pptx import Presentation
except ImportError:
    print(
        'Error: python-pptx is not installed. Run:\n'
        '  pnpm run setup-deck-python',
        file=sys.stderr,
    )
    sys.exit(1)


def normalize_note_text(text: str) -> str:
    if not text:
        return ''
    return ' '.join(text.split())


def main() -> None:
    parser = argparse.ArgumentParser(
        description='Extract speaker notes from a PowerPoint file to stdout (numbered lines).',
    )
    parser.add_argument(
        'pptx',
        help='Path to .pptx file',
    )
    parser.add_argument(
        '-o',
        '--output',
        help='Write to this file instead of stdout (UTF-8)',
    )
    args = parser.parse_args()

    prs = Presentation(args.pptx)
    lines = []
    for i, slide in enumerate(prs.slides, start=1):
        text = ''
        if slide.has_notes_slide:
            ns = slide.notes_slide
            if ns.notes_text_frame is not None:
                text = normalize_note_text(ns.notes_text_frame.text or '')
        lines.append(f'{i}: {text}')

    out = '\n'.join(lines) + '\n'
    if args.output:
        with open(args.output, 'w', encoding='utf-8') as f:
            f.write(out)
    else:
        sys.stdout.write(out)


if __name__ == '__main__':
    try:
        main()
    except Exception as e:
        print(f'Error: {e}', file=sys.stderr)
        sys.exit(1)
